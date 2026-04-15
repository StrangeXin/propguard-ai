#!/usr/bin/env python3.14
"""Generate AI Financial Trading Tools Market Prospect PPT"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Color Palette ──
BG_DARK    = RGBColor(0x0F, 0x17, 0x29)  # deep navy
BG_CARD    = RGBColor(0x16, 0x20, 0x3A)  # card bg
ACCENT     = RGBColor(0x38, 0xBD, 0xF8)  # bright cyan
ACCENT2    = RGBColor(0x22, 0xD3, 0xEE)  # teal
GREEN      = RGBColor(0x4A, 0xDE, 0x80)  # green
RED        = RGBColor(0xF8, 0x71, 0x71)  # red
ORANGE     = RGBColor(0xFB, 0xBF, 0x24)  # orange/yellow
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GRAY       = RGBColor(0x94, 0xA3, 0xB8)  # slate gray
LIGHT_GRAY = RGBColor(0xCB, 0xD5, 0xE1)

prs = Presentation()
prs.slide_width  = Inches(16)
prs.slide_height = Inches(9)

# ── Helper functions ──

def add_bg(slide):
    """Fill slide with dark background"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Arial"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multi_text(slide, left, top, width, height, lines, default_size=16):
    """lines = [(text, size, color, bold), ...]"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, size, color, bold) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = "Arial"
        p.space_after = Pt(6)
    return txBox

def add_rounded_rect(slide, left, top, width, height, fill_color=BG_CARD):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_section_header(slide, number, title):
    """Section divider style"""
    add_text_box(slide, 1.5, 2.5, 13, 1, f"PART {number}", 20, ACCENT, True)
    add_text_box(slide, 1.5, 3.2, 13, 2, title, 44, WHITE, True)

def add_page_number(slide, num, total=20):
    add_text_box(slide, 14.5, 8.3, 1.2, 0.4, f"{num}/{total}", 10, GRAY, alignment=PP_ALIGN.RIGHT)

def add_top_bar(slide, title, subtitle=""):
    # accent line
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(16), Inches(0.06))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    add_text_box(slide, 0.8, 0.3, 14, 0.6, title, 28, WHITE, True)
    if subtitle:
        add_text_box(slide, 0.8, 0.85, 14, 0.4, subtitle, 14, GRAY)

def make_table(slide, left, top, width, rows_data, col_widths=None):
    """rows_data = [["h1","h2",...], ["r1c1","r1c2",...], ...]"""
    rows = len(rows_data)
    cols = len(rows_data[0])
    table_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top),
                                          Inches(width), Inches(0.4 * rows))
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    for r, row in enumerate(rows_data):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.name = "Arial"
                if r == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                else:
                    paragraph.font.color.rgb = LIGHT_GRAY
            # cell fill
            cell.fill.solid()
            if r == 0:
                cell.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x4A)
            elif r % 2 == 0:
                cell.fill.fore_color.rgb = RGBColor(0x13, 0x1C, 0x34)
            else:
                cell.fill.fore_color.rgb = BG_CARD
    return table_shape


# ══════════════════════════════════════════════════════════════
# P1 - Cover
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)

# big accent block
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(16), Inches(0.12))
shape.fill.solid(); shape.fill.fore_color.rgb = ACCENT; shape.line.fill.background()

add_text_box(slide, 1.5, 2.0, 13, 1.2, "AI 金融交易工具", 52, WHITE, True)
add_text_box(slide, 1.5, 3.3, 13, 1.0, "市场前景与战略机遇", 40, ACCENT, True)

add_multi_text(slide, 1.5, 5.0, 13, 2, [
    ("基于近 30 天全网 420+ 条跨平台数据的深度调研", 18, GRAY, False),
    ("数据来源：Reddit / X / TikTok / Instagram / HN / GitHub / Web / Polymarket", 14, GRAY, False),
    ("", 10, GRAY, False),
    ("2026 年 4 月 13 日", 16, LIGHT_GRAY, False),
])
add_page_number(slide, 1)


# ══════════════════════════════════════════════════════════════
# P2 - Shocking Number
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_top_bar(slide, "市场正在发生什么")

add_text_box(slide, 1.5, 2.0, 13, 1.5, "85%", 120, ACCENT, True)
add_text_box(slide, 1.5, 4.2, 13, 1, "的外汇日交易量已由 AI 交易机器人驱动", 32, WHITE, True)
add_text_box(slide, 1.5, 5.5, 13, 0.6, "日均交易量 9.6 万亿美元  |  来源：@DataconomyMedia, 2026-04-12", 16, GRAY)

add_rounded_rect(slide, 1.5, 6.5, 6, 1.5)
add_multi_text(slide, 1.8, 6.6, 5.5, 1.3, [
    ("Claude AI Bot", 14, ACCENT, True),
    ("Polymarket: $1 → $330 万", 20, GREEN, True),
    ("来源: Hacker News / finbold.com", 11, GRAY, False),
])

add_rounded_rect(slide, 8.5, 6.5, 6, 1.5)
add_multi_text(slide, 8.8, 6.6, 5.5, 1.3, [
    ("中国程序员 AI Bot", 14, ACCENT, True),
    ("每 15 分钟数万笔 BTC 交易", 20, GREEN, True),
    ("来源: TikTok @finalrender_", 11, GRAY, False),
])
add_page_number(slide, 2)


# ══════════════════════════════════════════════════════════════
# P3 - Market Overview: 6 Tracks
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_top_bar(slide, "2026 年 AI 交易生态：六大热门赛道")

tracks = [
    ("自然语言策略生成", "极强", "Horizon AI / Robonet / Superior Trade", ACCENT),
    ("AI 自动化交易 Bot", "极强", "Pips Connect / AlgoAiden / Agent-Kai", ACCENT),
    ("AI 交易日记 & 分析", "强", "TradeDeck / TradesViz / LunarLog", GREEN),
    ("回测可视化工具", "中强", "MarketFlux / nautilus_trader", GREEN),
    ("跨平台信号复制", "中", "Telegram Signal Copier / ConnectXCopy", ORANGE),
    ("DeFi / 链上自动化", "增长中", "Hyperliquid 生态 / go-trader", ORANGE),
]

for i, (name, strength, players, color) in enumerate(tracks):
    row = i // 2
    col = i % 2
    x = 0.8 + col * 7.5
    y = 1.6 + row * 2.2

    add_rounded_rect(slide, x, y, 7, 1.9)
    add_text_box(slide, x + 0.3, y + 0.15, 4, 0.5, name, 18, WHITE, True)

    # strength badge
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(x + 5.2), Inches(y + 0.15), Inches(1.4), Inches(0.4))
    badge.fill.solid(); badge.fill.fore_color.rgb = color; badge.line.fill.background()
    badge.text_frame.paragraphs[0].text = strength
    badge.text_frame.paragraphs[0].font.size = Pt(11)
    badge.text_frame.paragraphs[0].font.color.rgb = BG_DARK
    badge.text_frame.paragraphs[0].font.bold = True
    badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_text_box(slide, x + 0.3, y + 0.7, 6.4, 0.5, f"代表产品：{players}", 12, GRAY)
add_page_number(slide, 3)


# ══════════════════════════════════════════════════════════════
# P4 - Section Divider: Part 2
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_header(slide, "02", "传统交易 vs AI 交易")
add_text_box(slide, 1.5, 5.5, 13, 1, "从手动盯盘到 AI 自主决策的范式转变", 20, GRAY)
add_page_number(slide, 4)


# ══════════════════════════════════════════════════════════════
# P5 - Traditional Pain Points
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_top_bar(slide, "传统交易平台的六大痛点")

pains = [
    ("平台割裂", "6+ 平台 API 互不相通\nMT4/MT5/cTrader/TradeLocker/DXTrade", "@connectxcopy"),
    ("编码门槛高", "策略自动化需学 MQL5/C#/Pine\n非程序员被拒之门外", "fortraders.com"),
    ("手动操作多", "盯盘、录日记、算仓位\n交易者疲于应付", "TradeDeck"),
    ("信息过载", "交易者花一半时间\n理解市场在发生什么", "@labtrade_"),
    ("情绪化决策", "恐惧与贪婪导致\n非理性交易行为", "innotechtoday.com"),
    ("学习曲线陡", "换个平台就不会用\n因此挂掉挑战账户", "@fortunemmxm"),
]

for i, (title, desc, src) in enumerate(pains):
    row = i // 3
    col = i % 3
    x = 0.6 + col * 5.1
    y = 1.5 + row * 3.5

    add_rounded_rect(slide, x, y, 4.7, 3.0)
    # number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.3), Inches(y + 0.3), Inches(0.5), Inches(0.5))
    circle.fill.solid(); circle.fill.fore_color.rgb = RED; circle.line.fill.background()
    circle.text_frame.paragraphs[0].text = str(i + 1)
    circle.text_frame.paragraphs[0].font.size = Pt(14)
    circle.text_frame.paragraphs[0].font.color.rgb = WHITE
    circle.text_frame.paragraphs[0].font.bold = True
    circle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_text_box(slide, x + 1.0, y + 0.3, 3.5, 0.5, title, 18, WHITE, True)
    add_text_box(slide, x + 0.3, y + 1.0, 4.1, 1.2, desc, 13, LIGHT_GRAY)
    add_text_box(slide, x + 0.3, y + 2.4, 4.1, 0.4, f"来源: {src}", 10, GRAY)
add_page_number(slide, 5)


# ══════════════════════════════════════════════════════════════
# P6 - Traditional vs AI Comparison Table
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_top_bar(slide, "AI 交易 vs 传统交易：六维度对比")

table_data = [
    ["维度", "传统方式", "AI 方式", "提升幅度"],
    ["策略创建", "学编程语言 (MQL5/C#)", "自然语言描述即可", "门槛降低 90%"],
    ["交易执行", "手动下单 / 简单 EA", "AI Agent 24/7 自主决策", "效率 x10"],
    ["风险管理", "手动设止损", "AI 动态仓位 + 自动风控", "纪律性 x5"],
    ["交易复盘", "手动录入日记", "截图识别 + AI 自动分析", "耗时减少 80%"],
    ["市场分析", "看图表、读新闻", "AI 盘前简报 + 实时信号", "覆盖面 x20"],
    ["跨平台管理", "分别登录各平台", "一个 API 聚合 20+ 经纪商", "统一管理"],
]
make_table(slide, 0.8, 1.5, 14.4, table_data, col_widths=[2.2, 4, 4, 2.8])
add_page_number(slide, 6)


# ══════════════════════════════════════════════════════════════
# P7 - Three Core Advantages
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_top_bar(slide, "AI 交易的三大核心优势")

advantages = [
    ("去编程化", "用自然语言描述策略\nAI 自动生成可执行逻辑", "Horizon AI / TradeLocker\nRobonet prompt-to-quant", ACCENT),
    ("去情绪化", "AI 不会恐惧也不会贪婪\n严格执行策略纪律", "多个 AI Bot 核心卖点\n减少人为非理性决策", GREEN),
    ("去人工化", "不到 1 小时构建 Bot\n24/7 全自动运行", "TikTok @markus864\nClaude Code 构建 Bot", ORANGE),
]

for i, (title, desc, example, color) in enumerate(advantages):
    x = 0.8 + i * 5.1
    y = 1.8

    add_rounded_rect(slide, x, y, 4.7, 5.5)

    # icon number
    num_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        Inches(x + 1.6), Inches(y + 0.4), Inches(1.5), Inches(1.5))
    num_shape.fill.solid(); num_shape.fill.fore_color.rgb = color; num_shape.line.fill.background()
    num_shape.text_frame.paragraphs[0].text = f"0{i+1}"
    num_shape.text_frame.paragraphs[0].font.size = Pt(36)
    num_shape.text_frame.paragraphs[0].font.color.rgb = BG_DARK
    num_shape.text_frame.paragraphs[0].font.bold = True
    num_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_text_box(slide, x + 0.3, y + 2.2, 4.1, 0.6, title, 24, WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide, x + 0.3, y + 2.9, 4.1, 1.0, desc, 14, LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + 0.3, y + 4.2, 4.1, 1.0, example, 12, GRAY, alignment=PP_ALIGN.CENTER)
add_page_number(slide, 7)


# ══════════════════════════════════════════════════════════════
# P8 - Section: Competitive Landscape
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_header(slide, "03", "竞争格局深度分析")
add_text_box(slide, 1.5, 5.5, 13, 1, "外汇经纪商 / 加密交易所 / AI 工具 / 开源生态", 20, GRAY)
add_page_number(slide, 8)


# ══════════════════════════════════════════════════════════════
# P9 - Forex Broker Landscape
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_top_bar(slide, "外汇赛道：经纪商 & 平台格局", "数据来源: ForexBrokers.com, 55brokers.com, AquaFutures")

# Platform comparison
table_data = [
    ["平台", "生态成熟度", "自动化语言", "自然语言策略", "Prop Firm 采用"],
    ["MT5", "最成熟，EA 商城丰富", "MQL5", "不支持", "广泛"],
    ["cTrader", "中等，cAlgo 生态", "C# / cAlgo", "不支持", "增长中"],
    ["TradeLocker", "新兴，快速增长", "Pine Script", "支持", "24 家已采用"],
]
make_table(slide, 0.8, 1.5, 14.4, table_data, col_widths=[2.2, 3.5, 2.8, 2.5, 3.0])

# Broker API
add_text_box(slide, 0.8, 4.2, 14, 0.5, "主流经纪商 API 能力", 18, ACCENT, True)
table_data2 = [
    ["经纪商", "API 类型", "特点", "适合场景"],
    ["Interactive Brokers", "TWS API / REST", "最全面，多资产，低延迟", "专业量化"],
    ["OANDA", "REST API v20", "文档优秀，易于集成", "Bot 开发入门首选"],
    ["Saxo Bank", "OpenAPI", "机构级，多资产类", "高端客户"],
    ["FxPro", "MT4/MT5/cTrader", "跨平台同一点差佣金", "灵活选择"],
]
make_table(slide, 0.8, 4.8, 14.4, table_data2, col_widths=[3, 3, 4.5, 3.5])
add_page_number(slide, 9)


# ══════════════════════════════════════════════════════════════
# P10 - Crypto Exchange Landscape
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_top_bar(slide, "加密货币赛道：交易所费率对比", "数据来源: ventureburn.com, @reboundx_net, plisio.net")

table_data = [
    ["交易所", "用户量", "Maker 费率", "Taker 费率", "返佣", "定位"],
    ["Binance", "2 亿+", "0.012%", "0.024%", "40%+10%", "全球最大，流动性最好"],
    ["OKX", "大量", "0.009%", "0.0225%", "55%+10%", "费率最低"],
    ["Bybit", "大量", "0.013%", "0.0286%", "35%+10%", "合约交易强"],
    ["Coinbase", "1 亿+", "较高", "较高", "-", "美国合规首选"],
    ["Kraken", "-", "中等", "中等", "-", "专业交易者，API 强"],
]
make_table(slide, 0.8, 1.4, 14.4, table_data, col_widths=[2, 1.5, 1.8, 1.8, 1.8, 4])

# Hyperliquid callout
add_rounded_rect(slide, 0.8, 5.5, 14.4, 2.8, RGBColor(0x14, 0x2A, 0x40))
add_text_box(slide, 1.2, 5.6, 6, 0.5, "Hyperliquid — DeFi 交易新星", 22, ACCENT, True)
add_multi_text(slide, 1.2, 6.2, 6.5, 2, [
    ("链上订单簿，接近 CEX 体验", 14, LIGHT_GRAY, False),
    ("无需 KYC，Agent Wallet 机制", 14, LIGHT_GRAY, False),
    ("支持多腿策略、跨资产风控", 14, LIGHT_GRAY, False),
    ("官方 Python SDK + 开源 Copy Bot", 14, LIGHT_GRAY, False),
])
add_multi_text(slide, 8.5, 6.2, 6, 2, [
    ("Polymarket 预测", 14, ORANGE, True),
    ("登陆 Binance 概率: 32%", 18, WHITE, True),
    ("", 6, GRAY, False),
    ("安全警告: MEV 三明治攻击每天数千次", 13, RED, False),
])
add_page_number(slide, 10)


# ══════════════════════════════════════════════════════════════
# P11 - AI Tool Competitive Matrix
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_top_bar(slide, "AI 交易工具竞品矩阵", "按 用户类型 x 功能方向 分类")

# 2x2 matrix
# draw axes
add_text_box(slide, 7.2, 1.3, 2, 0.4, "分析工具", 14, GRAY, True, PP_ALIGN.CENTER)
add_text_box(slide, 7.2, 8.0, 2, 0.4, "自动执行", 14, GRAY, True, PP_ALIGN.CENTER)
add_text_box(slide, 0.5, 4.5, 1.5, 0.4, "零代码\n用户", 12, GRAY, True, PP_ALIGN.CENTER)
add_text_box(slide, 14.5, 4.5, 1.5, 0.4, "开发者", 12, GRAY, True, PP_ALIGN.CENTER)

# quadrants
quads = [
    (1.5, 1.6, 6.3, 3.2, "零代码 + 分析", [
        "TradeDeck (AI 截图识别)",
        "TradesViz (AI 盘前分析)",
        "LunarLog / Profit App",
    ], GREEN),
    (8.5, 1.6, 6.3, 3.2, "开发者 + 分析", [
        "Kavout (AI 金融研究)",
        "Rafa.ai (投资组合 AI)",
        "TrendSpider (技术分析)",
    ], ACCENT),
    (1.5, 5.0, 6.3, 3.2, "零代码 + 执行", [
        "Pips Connect (2500+ Bot)",
        "AlgoAiden (自改进 Bot)",
        "BullGPT / Horizon AI",
        "竞争最激烈，无明确赢家",
    ], ORANGE),
    (8.5, 5.0, 6.3, 3.2, "开发者 + 执行", [
        "Robonet (prompt-to-quant)",
        "AI-Trader (开源全自动)",
        "Agent-Kai (自学习终端)",
    ], ACCENT2),
]

for x, y, w, h, title, items, color in quads:
    add_rounded_rect(slide, x, y, w, h)
    add_text_box(slide, x + 0.2, y + 0.15, w - 0.4, 0.4, title, 15, color, True)
    for j, item in enumerate(items):
        c = RED if "竞争最激烈" in item else LIGHT_GRAY
        b = True if "竞争最激烈" in item else False
        add_text_box(slide, x + 0.4, y + 0.6 + j * 0.55, w - 0.6, 0.4, f"  {item}", 12, c, b)
add_page_number(slide, 11)


# ══════════════════════════════════════════════════════════════
# P12 - Trading Journal Tools
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_top_bar(slide, "交易日记 & 分析工具竞品对比")

table_data = [
    ["产品", "核心亮点", "差异化", "来源"],
    ["TradeDeck", "AI 截图识别交易记录", "最创新的录入方式", "X @tradedeckapp"],
    ["TradesViz", "AI 盘前分析 + 持续升级 UI", "分析深度最强", "X @tradesviz"],
    ["TradeLens Pro", "AI 问答: \"哪个策略胜率最高?\"", "唯一平价 AI ($12.50/月)", "tradelens.vip"],
    ["Tradezella", "全功能一站式日记系统", "功能最全面", "TikTok 推广"],
    ["Traderwaves", "自动跟踪，免费起步", "最低使用门槛", "Trustpilot 好评"],
    ["LunarLog", "AI 交易日记 (新产品)", "TikTok 热推", "TikTok @moonboy_matt"],
    ["Rafa.ai", "AI 投资组合分析", "上传即分析", "X @fenrirnft"],
]
make_table(slide, 0.8, 1.4, 14.4, table_data, col_widths=[2.2, 3.8, 3.2, 3])

add_rounded_rect(slide, 0.8, 6.0, 14.4, 2, RGBColor(0x14, 0x2A, 0x40))
add_text_box(slide, 1.2, 6.1, 13, 0.5, "关键洞察", 16, ACCENT, True)
add_multi_text(slide, 1.2, 6.6, 13, 1.3, [
    ("\"你的交易日记是你最被低估的工具\" — @optiondrops (X)", 14, LIGHT_GRAY, False),
    ("Traderwaves 用户: \"every trader that hates manual journal input\" — Trustpilot", 14, LIGHT_GRAY, False),
    ("有交易者在用 AI 自建日记工具 — 说明现有产品仍有不足", 14, ORANGE, True),
])
add_page_number(slide, 12)


# ══════════════════════════════════════════════════════════════
# P13 - Developer Tools & Open Source
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_top_bar(slide, "开发者工具链 & 开源生态", "构建 AI 交易 Bot 的门槛已大幅降低")

table_data = [
    ["工具", "用途", "说明"],
    ["CCXT", "统一交易所 API", "100+ 交易所，一套代码适配"],
    ["TradingView Charts", "前端 K 线图表", "轻量级，开源"],
    ["Hyperliquid SDK", "DEX 接口", "官方 Python SDK"],
    ["AI-Trader (HKUDS)", "全自动 Agent 交易", "开源，兼容 Binance/Coinbase/IB"],
    ["Claude API", "AI 策略解析 / 分析", "自然语言理解 + 代码生成"],
    ["Connect Trade API", "统一经纪商接口", "一个 API 连接 20+ 合规经纪商"],
]
make_table(slide, 0.8, 1.4, 14.4, table_data, col_widths=[3, 3.5, 7.5])

add_rounded_rect(slide, 0.8, 5.8, 14.4, 2.5, RGBColor(0x14, 0x2A, 0x40))
add_text_box(slide, 1.2, 5.9, 13, 0.5, "实际构建案例", 16, ACCENT, True)
add_multi_text(slide, 1.2, 6.4, 6, 1.8, [
    ("Claude Code 交易 Bot — 不到 1 小时", 14, GREEN, True),
    ("@markus864: $200 压力测试，目标月增 3-10%", 12, LIGHT_GRAY, False),
    ("7 交易所套利扫描器 — 拒绝付 $29/月", 14, GREEN, True),
    ("DEV Community: 覆盖 3700 个 USDT 永续合约", 12, LIGHT_GRAY, False),
])
add_multi_text(slide, 8, 6.4, 6, 1.8, [
    ("Hyperliquid Perp Bot + TUI Dashboard", 14, GREEN, True),
    ("GitHub: 完整均值回归策略实现", 12, LIGHT_GRAY, False),
    ("自定义算法交易 + Angel One 直连", 14, GREEN, True),
    ("Instagram: 无第三方依赖的自建图表", 12, LIGHT_GRAY, False),
])
add_page_number(slide, 13)


# ══════════════════════════════════════════════════════════════
# P14 - Section: Opportunities & Risks
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_header(slide, "04", "机遇与风险")
add_text_box(slide, 1.5, 5.5, 13, 1, "五个切入方向 / 市场亮点 / 核心挑战", 20, GRAY)
add_page_number(slide, 14)


# ══════════════════════════════════════════════════════════════
# P15 - Five Directions
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_top_bar(slide, "五个可切入的方向")

table_data = [
    ["方向", "开发周期", "预计 MRR", "技术难度", "推荐指数"],
    ["自然语言策略生成器", "4-6 周", "$29-99/用户", "中高", "最大机会 ★★★★★"],
    ["AI 交易日记 (截图识别)", "2-3 周", "$15-30/用户", "中", "最快启动 ★★★★"],
    ["跨平台交易聚合器", "6-8 周", "$39/用户", "高", "壁垒最深 ★★★"],
    ["Telegram 信号 Bot", "1-2 周", "$19-49/用户", "低", "最快变现 ★★★★"],
    ["AI 量化回测可视化", "3-4 周", "$25/用户", "中", "技术型 ★★★"],
]
make_table(slide, 0.8, 1.4, 14.4, table_data, col_widths=[3.5, 2, 2.5, 2, 4])

add_rounded_rect(slide, 0.8, 5.5, 7, 2.8, RGBColor(0x14, 0x2A, 0x40))
add_multi_text(slide, 1.2, 5.6, 6, 2.6, [
    ("推荐首选: Telegram 信号 Bot", 16, GREEN, True),
    ("", 4, GRAY, False),
    ("开发最快 (1-2 周)", 14, LIGHT_GRAY, False),
    ("需求明确，付费意愿强", 14, LIGHT_GRAY, False),
    ("法律风险低 (工具定位)", 14, LIGHT_GRAY, False),
    ("Telegram Signal Copier 已验证市场: 90,000+ 用户", 13, GRAY, False),
])

add_rounded_rect(slide, 8.3, 5.5, 7, 2.8, RGBColor(0x14, 0x2A, 0x40))
add_multi_text(slide, 8.7, 5.6, 6, 2.6, [
    ("最大机会: 自然语言策略生成", 16, ACCENT, True),
    ("", 4, GRAY, False),
    ("市场天花板最高", 14, LIGHT_GRAY, False),
    ("多个竞品但无明确赢家", 14, LIGHT_GRAY, False),
    ("需要更多打磨时间 (4-6 周)", 14, LIGHT_GRAY, False),
    ("Horizon AI / Robonet / Superior Trade 均在赛道中", 13, GRAY, False),
])
add_page_number(slide, 15)


# ══════════════════════════════════════════════════════════════
# P16 - Why Now
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_top_bar(slide, "为什么现在是最佳时机")

highlights = [
    ("AI 编码工具爆发", "Vibe coding 使 App Store\n提交量增长 30%", ACCENT),
    ("成功先例已有", "Solo founder 29 天\n$953 收入，$0 营销", GREEN),
    ("基础设施成熟", "CCXT + Claude API +\nSupabase 一站式就绪", ACCENT2),
    ("竞争窗口期", "多个赛道有玩家\n但无明确赢家", ORANGE),
    ("付费意愿强", "交易者习惯为工具付费\n$12.50-$99/月", GREEN),
]

for i, (title, desc, color) in enumerate(highlights):
    x = 0.6 + i * 3.05
    y = 2.0
    add_rounded_rect(slide, x, y, 2.8, 4.5)

    num_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                        Inches(x + 0.9), Inches(y + 0.4), Inches(1), Inches(1))
    num_shape.fill.solid(); num_shape.fill.fore_color.rgb = color; num_shape.line.fill.background()
    num_shape.text_frame.paragraphs[0].text = str(i + 1)
    num_shape.text_frame.paragraphs[0].font.size = Pt(28)
    num_shape.text_frame.paragraphs[0].font.color.rgb = BG_DARK
    num_shape.text_frame.paragraphs[0].font.bold = True
    num_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_text_box(slide, x + 0.15, y + 1.7, 2.5, 0.7, title, 16, WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide, x + 0.15, y + 2.6, 2.5, 1.5, desc, 13, LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_page_number(slide, 16)


# ══════════════════════════════════════════════════════════════
# P17 - Risks & Challenges
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_top_bar(slide, "风险与挑战")

risks = [
    ("法律合规", "佛罗里达州调查 OpenAI\nAI + 金融监管趋严", "需标注\"非投资建议\"", RED),
    ("技术维护", "Vibe coding 第 3 个月\n面临维护困难", "MVP 验证后重构核心", ORANGE),
    ("虚假宣传", "\"月收益 10%\" 等不实承诺\n与骗局产品混淆", "以工具定位拉开差距", ORANGE),
    ("MEV 安全", "DeFi 三明治攻击\n每天发生数千次", "链上交易需防范", RED),
    ("API 依赖", "交易所/经纪商 API\n可能随时变更限流", "多源冗余设计", ORANGE),
    ("竞争加速", "市场热度高\n新进入者不断涌入", "速度是最好的壁垒", ORANGE),
]

for i, (title, desc, solution, color) in enumerate(risks):
    row = i // 3
    col = i % 3
    x = 0.6 + col * 5.1
    y = 1.5 + row * 3.7

    add_rounded_rect(slide, x, y, 4.7, 3.2)
    # warning icon
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(x + 0.3), Inches(y + 0.25), Inches(1.2), Inches(0.35))
    badge.fill.solid(); badge.fill.fore_color.rgb = color; badge.line.fill.background()
    badge.text_frame.paragraphs[0].text = title
    badge.text_frame.paragraphs[0].font.size = Pt(11)
    badge.text_frame.paragraphs[0].font.color.rgb = BG_DARK
    badge.text_frame.paragraphs[0].font.bold = True
    badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_text_box(slide, x + 0.3, y + 0.8, 4.1, 1.0, desc, 13, LIGHT_GRAY)
    add_text_box(slide, x + 0.3, y + 2.2, 4.1, 0.5, f"应对: {solution}", 12, GREEN)
add_page_number(slide, 17)


# ══════════════════════════════════════════════════════════════
# P18 - Action Plan
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_top_bar(slide, "推荐行动方案：两步走策略")

# Step 1
add_rounded_rect(slide, 0.8, 1.5, 7, 6.5, RGBColor(0x14, 0x2A, 0x40))
step1_badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(1.2), Inches(1.7), Inches(2.5), Inches(0.5))
step1_badge.fill.solid(); step1_badge.fill.fore_color.rgb = GREEN; step1_badge.line.fill.background()
step1_badge.text_frame.paragraphs[0].text = "第一步: 快速 MVP"
step1_badge.text_frame.paragraphs[0].font.size = Pt(14)
step1_badge.text_frame.paragraphs[0].font.color.rgb = BG_DARK
step1_badge.text_frame.paragraphs[0].font.bold = True
step1_badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

add_text_box(slide, 4.2, 1.7, 3, 0.5, "0 - 4 周", 16, GREEN, True)

add_multi_text(slide, 1.2, 2.5, 6.2, 5, [
    ("产品选择", 15, ACCENT, True),
    ("Telegram 信号 Bot 或 AI 交易日记", 13, LIGHT_GRAY, False),
    ("", 8, GRAY, False),
    ("技术栈", 15, ACCENT, True),
    ("Claude Code + Supabase + Vercel", 13, LIGHT_GRAY, False),
    ("Python (FastAPI) + Stripe", 13, LIGHT_GRAY, False),
    ("", 8, GRAY, False),
    ("里程碑", 15, ACCENT, True),
    ("第 1 周: 框架搭建 + 核心功能原型", 13, LIGHT_GRAY, False),
    ("第 2 周: AI 集成 + 基础 UI + 内测", 13, LIGHT_GRAY, False),
    ("第 3 周: 支付集成 + Landing Page", 13, LIGHT_GRAY, False),
    ("第 4 周: Product Hunt 发布", 13, LIGHT_GRAY, False),
    ("", 8, GRAY, False),
    ("目标: 验证需求，获取前 100 个用户", 14, GREEN, True),
])

# Step 2
add_rounded_rect(slide, 8.3, 1.5, 7, 6.5, RGBColor(0x14, 0x2A, 0x40))
step2_badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(8.7), Inches(1.7), Inches(2.8), Inches(0.5))
step2_badge.fill.solid(); step2_badge.fill.fore_color.rgb = ACCENT; step2_badge.line.fill.background()
step2_badge.text_frame.paragraphs[0].text = "第二步: 核心产品"
step2_badge.text_frame.paragraphs[0].font.size = Pt(14)
step2_badge.text_frame.paragraphs[0].font.color.rgb = BG_DARK
step2_badge.text_frame.paragraphs[0].font.bold = True
step2_badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

add_text_box(slide, 12, 1.7, 3, 0.5, "1 - 3 个月", 16, ACCENT, True)

add_multi_text(slide, 8.7, 2.5, 6.2, 5, [
    ("产品升级", 15, ACCENT, True),
    ("转向自然语言策略生成器", 13, LIGHT_GRAY, False),
    ("复用 MVP 阶段用户反馈和技术积累", 13, LIGHT_GRAY, False),
    ("", 8, GRAY, False),
    ("营销渠道", 15, ACCENT, True),
    ("1. Reddit (r/Daytrading, r/algotrading)", 13, LIGHT_GRAY, False),
    ("2. TikTok (AI 交易内容热度最高)", 13, LIGHT_GRAY, False),
    ("3. X/Twitter (金融科技创作者)", 13, LIGHT_GRAY, False),
    ("4. Product Hunt (首发)", 13, LIGHT_GRAY, False),
    ("", 8, GRAY, False),
    ("目标: $5K MRR", 14, ACCENT, True),
])
add_page_number(slide, 18)


# ══════════════════════════════════════════════════════════════
# P19 - Three Core Judgments (Summary)
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_top_bar(slide, "总结：三个核心判断")

judgments = [
    ("01", "AI 交易不是未来\n是现在", "85% 外汇交易量\n已由 AI 驱动", ACCENT),
    ("02", "工具型产品\n最安全", "做工具不做决策\n规避法律风险", GREEN),
    ("03", "窗口期有限\n2026 是最佳进入时间", "等待即是被淘汰\n速度是最好的壁垒", ORANGE),
]

for i, (num, title, desc, color) in enumerate(judgments):
    x = 0.8 + i * 5.1
    y = 2.0

    add_rounded_rect(slide, x, y, 4.7, 5.5)

    num_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        Inches(x + 1.2), Inches(y + 0.5), Inches(2.3), Inches(1.5))
    num_shape.fill.solid(); num_shape.fill.fore_color.rgb = color; num_shape.line.fill.background()
    num_shape.text_frame.paragraphs[0].text = num
    num_shape.text_frame.paragraphs[0].font.size = Pt(48)
    num_shape.text_frame.paragraphs[0].font.color.rgb = BG_DARK
    num_shape.text_frame.paragraphs[0].font.bold = True
    num_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_text_box(slide, x + 0.3, y + 2.3, 4.1, 1.2, title, 20, WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide, x + 0.3, y + 3.8, 4.1, 1.2, desc, 14, LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_page_number(slide, 19)


# ══════════════════════════════════════════════════════════════
# P20 - Appendix / Data Sources
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_top_bar(slide, "附录：数据来源与方法论")

add_multi_text(slide, 0.8, 1.4, 7, 3.5, [
    ("调研工具", 16, ACCENT, True),
    ("last30days v3.0.0 跨平台数据聚合引擎", 14, LIGHT_GRAY, False),
    ("", 8, GRAY, False),
    ("数据范围", 16, ACCENT, True),
    ("2026-03-14 至 2026-04-13（近 30 天）", 14, LIGHT_GRAY, False),
    ("", 8, GRAY, False),
    ("总证据量", 16, ACCENT, True),
    ("420+ 条原始证据，经 AI 去重、排序、聚类", 14, LIGHT_GRAY, False),
    ("", 8, GRAY, False),
    ("搜索轮次", 16, ACCENT, True),
    ("8 轮并行搜索，覆盖 4 个维度", 14, LIGHT_GRAY, False),
])

add_multi_text(slide, 8.5, 1.4, 7, 3.5, [
    ("覆盖平台", 16, ACCENT, True),
    ("Reddit      — 社区讨论与用户反馈", 13, LIGHT_GRAY, False),
    ("X/Twitter   — 实时行业动态", 13, LIGHT_GRAY, False),
    ("TikTok      — 创作者生态 (热度最高)", 13, LIGHT_GRAY, False),
    ("Instagram   — 视觉化内容", 13, LIGHT_GRAY, False),
    ("GitHub      — 开源项目与技术趋势", 13, LIGHT_GRAY, False),
    ("Brave Web   — 专业文章与评测", 13, LIGHT_GRAY, False),
    ("Hacker News — 技术社区观点", 13, LIGHT_GRAY, False),
    ("Polymarket  — 预测市场信号", 13, LIGHT_GRAY, False),
])

add_text_box(slide, 0.8, 5.5, 14.4, 0.5, "关键参考链接 (Top 10)", 16, ACCENT, True)

refs = [
    ["来源", "标题", "链接"],
    ["@DataconomyMedia", "AI bots 驱动 85% 外汇交易", "x.com/DataconomyMedia/status/2043207538875802067"],
    ["ForexBrokers.com", "7 Best Forex Trading APIs 2026", "forexbrokers.com/guides/best-api-brokers"],
    ["ventureburn.com", "12 Crypto Exchanges Lowest Fees", "ventureburn.com/crypto-exchange-with-lowest-fees/"],
    ["hiperwire.io", "AI Trading Bots on Hyperliquid", "hiperwire.io/explainers/ai-trading-bots-agents-hyperliquid"],
    ["fortraders.com", "AI Trading Tools That Work", "fortraders.com/blog/ai-trading-tools-work"],
    ["@RoundtableSpace", "Trading GitHub 开发者资源清单", "x.com/RoundtableSpace/status/2043369805495906756"],
    ["finbold.com", "Claude Bot: $1 → $3.3M", "finbold.com/claude-ai-powered-trading-bot-turns-1-into-3-3-million"],
    ["DEV Community", "7 交易所套利扫描器", "dev.to/foxyyybusiness/...funding-rate-arbitrage-scanner"],
    ["tradelens.vip", "Best Trade Journal Apps 对比", "tradelens.vip/resources/best-trade-journal-apps"],
    ["Fortune", "Cursor CEO: Vibe Coding 警告", "fortune.com/article/cursor-ceo-vibe-coding-warning"],
]
make_table(slide, 0.8, 6.0, 14.4, refs, col_widths=[3, 4.5, 6.5])
add_page_number(slide, 20)


# ── Save ──
output_path = os.path.join(os.path.dirname(__file__), "AI金融交易工具_市场前景汇报.pptx")
prs.save(output_path)
print(f"PPT saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
